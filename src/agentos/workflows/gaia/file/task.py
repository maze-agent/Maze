import os
import ray
import time
import gc
import torch
import json
import pdfplumber  #  PDF handle
import pandas as pd
from io import BytesIO
from docx import Document
from pptx import Presentation
import xml.etree.ElementTree as ET
from xml.dom import minidom
from typing import List, Dict, Tuple

from agentos.scheduler import cpu, gpu, io
from agentos.utils.remote_llm_route import RemoteLlmRoute
import cloudpickle
import re
def estimate_tokens(text):
    """
    A slightly more accurate tokenizer estimator for mixed Chinese/English text.
    It counts CJK characters and non-CJK tokens separately to avoid double counting.
    """
    # 1. CJK
    cjk_chars = sum(1 for char in text if '\u4E00' <= char <= '\u9FFF')
    # 2. CJK
    # CJK
    non_cjk_text = re.sub(r'[\u4E00-\u9FFF]', ' ', text)
    non_cjk_text = non_cjk_text.replace("\n", " ")
    # 3. CJK
    # usesplit()
    non_cjk_words_count = len(non_cjk_text.split())
    # 4. 
    # 1.3token
    estimated_tokens = cjk_chars + int(non_cjk_words_count * 1.3)
    return estimated_tokens

def query_vllm_model(
    endpoint: RemoteLlmRoute,
    model_alias: str,
    messages: List,
    temperature: float = 0.6,
    max_token: int = 1024,
    top_p: float = 0.9,
    repetition_penalty: float = 1.1,
) -> tuple[dict, str]:
    """Call a colocated vLLM worker via its OpenAI-compatible ``/v1/chat/completions`` route."""
    import requests
    from rich.console import Console
    console = Console()
    chat_url = endpoint.completions_url
    headers = endpoint.headers()
    # Build OpenAI-compatible payload
    payload = {
        "model": model_alias,
        "messages": messages,
        "temperature": temperature,
        "max_tokens": max_token,
        "top_p": top_p,
        "repetition_penalty": repetition_penalty,
    }
    #

    conversation = ""
    for message in messages:
        conversation += f"{message['role']}: {message['content']}"
    try:
        response = requests.post(chat_url, json=payload, headers=headers, timeout=3600)
        response.raise_for_status()  # raise_for_status on failure
        
        response_data = response.json()
        content = response_data['choices'][0]['message']['content'].lstrip()
        
        # queryaligned
        features = {"text_length": len(conversation), "token_count": estimate_tokens(conversation)}
        return features, content

    except requests.exceptions.RequestException as e:
        error_msg = f"vLLM request failed: {str(e)}"
        console.print(f"[bold red]{error_msg}")
        features = {"text_length": len(conversation), "token_count": estimate_tokens(conversation)}
        return features, f"[bold red]{error_msg}"

def query_llm_online(endpoint: RemoteLlmRoute, payload: Dict[str, str], tokenizer_path: str) -> Tuple[dict, str]:
    """POST to a hosted OpenAI-compatible chat endpoint (credentials live on the route object)."""
    import requests
    from rich.console import Console
    conversation = ""
    for message in payload["messages"]:
        conversation += f"{message['role']}: {message['content']}"

    console = Console()
    headers = endpoint.headers()
    try:
        response = requests.post(
            endpoint.completions_url,
            json=payload,
            headers=headers,
            timeout=3600,
        )
        response.raise_for_status()
        return {"text_length": len(conversation), "token_count": estimate_tokens(conversation)}, response.json()['choices'][0]['message']['content'].lstrip()
    except Exception as e:
        console.print(f"[bold red]API request failed: {str(e)}")
        return {"text_length": len(conversation), "token_count": estimate_tokens(conversation)}, f"[bold red]API request failed: {str(e)}"

def query_llm(model:str, model_folder: str, messages:List, temperature:float= 0.6, max_token:int= 1024, top_p:float= 0.9, repetition_penalty:float= 1.1):
    from transformers import AutoTokenizer, AutoModelForCausalLM
    import torch
    # Load tokenizer and weights
    tokenizer_path= os.path.join(model_folder, "Qwen/Qwen3-32B")
    tokenizer = AutoTokenizer.from_pretrained(tokenizer_path, device_map= "cuda")
    local_model = AutoModelForCausalLM.from_pretrained(
        os.path.join(model_folder, model),
        torch_dtype="float16",
        low_cpu_mem_usage=True,
        device_map="cuda", offload_state_dict= False,
        )

    # Flatten messages to a string
    conversation = ""
    for message in messages:
        conversation += f"{message['role']}: {message['content']}"

    # Tokenize prompt
    input_ids = tokenizer.encode(conversation + tokenizer.eos_token, return_tensors='pt').to("cuda")

    # Generate completion
    output = local_model.generate(
        input_ids, 
        pad_token_id=tokenizer.eos_token_id,
        max_new_tokens= max_token,
        temperature= temperature, 
        top_p= top_p,
        repetition_penalty= repetition_penalty
    )
    response = tokenizer.decode(output[:, input_ids.shape[-1]:][0], skip_special_tokens=True)

    del local_model
    del tokenizer
    gc.collect()
    torch.cuda.empty_cache()
    return {"text_length": len(conversation), "token_count": estimate_tokens(conversation)}, response

@cpu(cpu_num= 1, mem= 1024)
def task1_file_process(context):
    """
    Processes the file content based on type and stores the processed content in context.
    """
    try:
        time_record = {
            "get_time": 0.0,
            "put_size_bytes": 0,
            "get_size_bytes": 0
        }
        start_time= time.time()
        dag_id= ray.get(context.get.remote('dag_id'))
        question = ray.get(context.get.remote("question"))
        supplementary_files = ray.get(context.get.remote("supplementary_files"))
        time_record["get_time"] = time.time()- start_time
        time_record["get_size_bytes"] = sum(
            len(cloudpickle.dumps(obj)) for obj in 
            [dag_id, question, supplementary_files]
        )
        file_name = ""
        content = None
        
        if supplementary_files:
            file_name = next(iter(supplementary_files.keys()), "")
            content = next(iter(supplementary_files.values()), None)
            print(f"  -> Found single supplementary file to process: '{file_name}'")
        else:
            print("  -> No supplementary files found in context.")
        processed_content = ""
        
        print("✅ Task 1: Starting file processing...")
        if content is None:
            print(f"  -> No content for file '{file_name}', skipping processing.")
        else:
            print(f"  -> Processing file: {file_name}")
            file_name_lower = file_name.lower()
            
            # Excel filehandle
            if file_name_lower.endswith(('.xlsx', '.xls')):
                print("handle Excel file...")
                excel_file = BytesIO(content)
                
                # sheet
                excel_data = []
                with pd.ExcelFile(excel_file) as xls:
                    for sheet_name in xls.sheet_names:
                        df = pd.read_excel(xls, sheet_name)
                        excel_data.append(f"=== Sheet: {sheet_name} ===\n")
                        excel_data.append(df.to_string(index=True))
                
                processed_content = "\n\n".join(excel_data)
                print(f"extracted to {len(excel_data)//2} content")

            # CSV filehandle
            elif file_name_lower.endswith('.csv'):
                print("handle CSV file...")
                csv_file = BytesIO(content)
                df = pd.read_csv(csv_file)
                processed_content = df.to_string(index=True)
                print(f"extracted to CSV contenttotal {len(df)} ")

            # XML filehandle
            elif file_name_lower.endswith('.xml'):
                print("[File Process] handle XML file...")
                try:
                    xml_content = content.decode('utf-8')
                    print("[File Process] XML content")
                    
                    #  XML
                    try:
                        parsed_xml = minidom.parseString(xml_content)
                        processed_content = parsed_xml.toprettyxml()
                        print("[File Process] XML ")
                    except Exception as xml_error:
                        print(f"[File Process] XML failed: {str(xml_error)}")
                        processed_content = xml_content
                    
                    #  XML 
                    try:
                        root = ET.fromstring(xml_content)
                        def extract_text(element):
                            text = element.text or ''
                            for child in element:
                                text += extract_text(child)
                                if child.tail:
                                    text += child.tail
                            return text
                        
                        text_content = extract_text(root)
                        print("[File Process] content")
                        
                        processed_content = (
                            "=== XML Structure ===\n"
                            f"{processed_content}\n\n"
                            "=== Extracted Text Content ===\n"
                            f"{text_content}"
                        )
                    except Exception as text_error:
                        print(f"[File Process] failed: {str(text_error)}")
                    
                    print(f"[File Process] XML handledonecontent: {len(processed_content)}")
                    
                except Exception as e:
                    print(f"[File Process] XML handle: {str(e)}")
                    processed_content = content.decode('utf-8')

            # PPT filehandle
            elif file_name_lower.endswith('.pptx'):
                print("handle PPT file...")
                prs = Presentation(BytesIO(content))
                all_slides_text = []
                for i, slide in enumerate(prs.slides):
                    slide_text_parts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text_parts.append(shape.text.strip())
                    all_slides_text.append(f"--- Slide {i+1} ---\n" + '\n'.join(filter(None, slide_text_parts)))
                processed_content = "\n\n".join(all_slides_text)
                print(f"extracted to {len(all_slides_text)} content")

            # Word filehandle
            elif file_name_lower.endswith('.docx'):
                print("handle Word file...")
                doc = Document(BytesIO(content))
                paragraphs_text = []
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if text:
                        paragraphs_text.append(text)
                processed_content = "\n".join(paragraphs_text)
                print(f"extracted to {len(paragraphs_text)} ")

            # PDF filehandle
            elif file_name_lower.endswith('.pdf'):
                print("handle PDF file...")
                temp_pdf_path = f"/tmp/temp_{file_name_lower.split('.')[0]}.pdf"
                with open(temp_pdf_path, "wb") as f:
                    f.write(content)
                
                pdf_pages = []
                with pdfplumber.open(temp_pdf_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            pdf_pages.append(f"--- Page {i+1} ---\n{page_text}")
                        
                        tables = page.extract_tables()
                        if tables:
                            table_texts = []
                            for j, table in enumerate(tables):
                                table_text = f"Table {j+1} on Page {i+1}:\n"
                                for row in table:
                                    cleaned_row = [str(cell) if cell is not None else "" for cell in row]
                                    table_text += " | ".join(cleaned_row) + "\n"
                                table_texts.append(table_text)
                            
                            if table_texts:
                                pdf_pages.append("\n".join(table_texts))
                
                os.remove(temp_pdf_path)
                processed_content = "\n\n".join(pdf_pages)
                print(f"extracted to {len(pdf_pages)} / PDF content")

            # filehandle
            elif file_name_lower.endswith('.txt'):
                print("handle TXT file...")
                processed_content = content.decode('utf-8')

            else:
                print(f"warning: file: {file_name_lower}handle")
                try:
                    processed_content = content.decode('utf-8')
                except UnicodeDecodeError:
                    print("file UTF-8content")
                    processed_content = "[Content could not be decoded]"
            print("✅ Task 1: File processing finished.")
        # model
        prompt= (
            "#Background#\n"
            "You are a general AI assistant. I will ask you a question. Report your thoughts, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER].\n"
            "YOUR FINAL ANSWER should be a number OR as few words as possible OR a comma separated list of numbers and/or strings.\n"
            "If you are asked for a number, don’t use comma to write your number neither use units such as $ or percent sign unless specified otherwise.\n"
            "If you are asked for a string, don’t use articles, neither abbreviations (e.g. for cities), and write the digits in plain text unless specified otherwise.\n"
            "If you are asked for a comma separated list, apply the above rules depending of whether the element to be put in the list is a number or a string.\n"
            f"#Question#\n{question}\n"
            "#Extracted text from file#\n"
            f"{processed_content[:8000]}\n"
        )
        llm_process_feat= {"text_length": len(prompt), "token_count": estimate_tokens(prompt)}
        time_record["put_size_bytes"] += len(cloudpickle.dumps(processed_content))
        # handle
        context.put.remote("processed_content", processed_content)
        print(f"✅ filehandle...")
        return json.dumps({ #

            "dag_id": dag_id,
            "succ_task_feat": {
                "task2_llm_process_qwen": {"text_length": llm_process_feat["text_length"], "token_count": llm_process_feat["token_count"], "reason": 1},
                "task3_llm_process_deepseek": {"text_length": llm_process_feat["text_length"], "token_count": llm_process_feat["token_count"], "reason": 0}
            },
            "curr_task_feat": None,
            "start_time": start_time,
            "end_time": time.time(),
            "time_record": time_record
        })
    except Exception as e:
        print(f"task2_file_process error: {str(e)}")
        raise e

@gpu(gpu_mem= 70000, model_name= "qwen3-32b", backend="huggingface")
def task2_llm_process_qwen(context):
    
    """
    Processes the file content using LLM based on the question.
    """
    try:
        time_record= {
            "get_time": 0.0,
            "put_size_bytes": 0,
            "get_size_bytes": 0
        }
        backend= task2_llm_process_qwen._task_decorator['backend']
        print(f"✅ qwenhandlestart....")
        start_time= time.time()  # ADD:  
        dag_id= ray.get(context.get.remote("dag_id"))
        processed_content = ray.get(context.get.remote("processed_content")) 
        question= ray.get(context.get.remote("question"))

        use_online_model= ray.get(context.get.remote("use_online_model"))
        model_folder= ray.get(context.get.remote("model_folder"))
        tokenizer_path= os.path.join(model_folder, "Qwen/Qwen3-32B")
        temperature, max_token, repetition_penalty, top_p= None, None, None, None
        if not use_online_model:
            temperature= ray.get(context.get.remote("temperature"))
            max_token= ray.get(context.get.remote("max_tokens"))
            top_p= ray.get(context.get.remote("top_p"))
            repetition_penalty= ray.get(context.get.remote("repetition_penalty"))
        if not question:
             raise ValueError(f"task {dag_id} missing Question")
        time_record["get_time"] = time.time() - start_time
        # Build prompt
        # model
        prompt= (
            "#Background#\n"
            "You are a general AI assistant. I will ask you a question. Report your thoughts, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER].\n"
            "YOUR FINAL ANSWER should be a number OR as few words as possible OR a comma separated list of numbers and/or strings.\n"
            "If you are asked for a number, don’t use comma to write your number neither use units such as $ or percent sign unless specified otherwise.\n"
            "If you are asked for a string, don’t use articles, neither abbreviations (e.g. for cities), and write the digits in plain text unless specified otherwise.\n"
            "If you are asked for a comma separated list, apply the above rules depending of whether the element to be put in the list is a number or a string.\n"
            f"#Question#\n{question}\n"
            "#Extracted text from file#\n"
            f"{processed_content[:8000]}\n"
        )
        payload= {
            "model": "Qwen/Qwen3-32B",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": temperature,
            "max_tokens": max_token,
            "enable_thinking": False,
            "response_format": {"type": "text"}               
        }
        inference_features, answer= None, None
        if use_online_model:
            inference_features, answer = query_llm_online(
                RemoteLlmRoute.from_dag_context(context),
                payload=payload,
                tokenizer_path=tokenizer_path,
            )
        elif backend == "vllm":
            inference_features, answer = query_vllm_model(
                RemoteLlmRoute.for_vllm_worker_base(
                    ray.get(context.get.remote(f"task2_llm_process_qwen_request_api_url"))
                ),
                model_alias="qwen3-32b",
                messages= payload["messages"],
                temperature= temperature,
                max_token= max_token,
                top_p= top_p,
                repetition_penalty= repetition_penalty)
        else:
            inference_features, answer= query_llm(model_folder= model_folder, model= "Qwen/Qwen3-32B", messages= [{"role": "user", "content": prompt}], temperature= temperature, max_token= max_token, top_p= top_p, repetition_penalty= repetition_penalty)

        inference_features['reason'] = 0
        context.put.remote("qwen_answer", answer)
        next_task_prompt = (
            "You are a senior editor and a world-class reasoning expert. Your job is to synthesize the answers from two different AI assistants to produce one final, superior answer for the given question.\n\n"
            f"--- Original Question ---\n{question}\n\n"
            f"--- Answer from Assistant 1 (Qwen3) ---"
            f"--- Answer from Assistant 2 (DeepSeek) ---"
            "--- Your Task ---\n"
            "Analyze both answers. Identify the strengths and weaknesses of each. Then, combine their best elements, correct any errors, and provide a single, comprehensive, and accurate final answer. Adhere to the final answer format requested in the original prompt.\n\n"
            "Report your thoughts, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER]."
        )
        text1_length= len(answer)
        text1_token_count= estimate_tokens(answer)
        text1_feature= {"text1_length": text1_length, "text1_token_count": text1_token_count}
        context.put.remote("text1_feature", text1_feature)
        time_record["get_size_bytes"]+= sum(
            len(cloudpickle.dumps(obj)) for obj in 
            [dag_id, processed_content, question, use_online_model, model_folder, 
             temperature, max_token, top_p, repetition_penalty]
        )
        time_record["put_size_bytes"]+= sum(
            len(cloudpickle.dumps(obj)) for obj in 
            [text1_feature]
        )
        return json.dumps({
            "dag_id": dag_id,
            "curr_task_feat": inference_features,
            "succ_task_feat": {
                "task4_llm_fuse_answer":{
                    "prompt_length": len(next_task_prompt),
                    "prompt_token_count": estimate_tokens(next_task_prompt),
                    "text1_length": text1_length,
                    "text1_token_count": text1_token_count,
                    "reason": 0
                }
            },
            "start_time": start_time,
            "end_time": time.time(),
            "time_record": time_record
        })
    except Exception as e:
        print(f"task2_llm_process_qwen error: {str(e)}")
        raise e

@gpu(gpu_mem= 70000, model_name= "deepseek-r1-32b", backend="huggingface")
def task3_llm_process_deepseek(context):
    
    """
    Processes the file content using LLM based on the question.
    """
    try:
        time_record = {
            "get_time": 0.0,
            "put_size_bytes": 0,
            "get_size_bytes": 0
        }
        backend= task3_llm_process_deepseek._task_decorator['backend']
        start_time= time.time()  # ADD:  
        print(f"✅ deepseekhandlestart....")
        dag_id= ray.get(context.get.remote("dag_id"))
        processed_content = ray.get(context.get.remote("processed_content")) 
        question= ray.get(context.get.remote("question"))

        use_online_model= ray.get(context.get.remote("use_online_model"))
        model_folder= ray.get(context.get.remote("model_folder"))
        tokenizer_path= os.path.join(model_folder, "Qwen/Qwen3-32B")
        temperature, max_token, repetition_penalty, top_p= None, None, None, None
        if not use_online_model:
            temperature= ray.get(context.get.remote("temperature"))
            max_token= ray.get(context.get.remote("max_tokens"))
            top_p= ray.get(context.get.remote("top_p"))
            repetition_penalty= ray.get(context.get.remote("repetition_penalty"))
        if not question:
             raise ValueError(f"task {dag_id} missing Question")
        time_record["get_time"] = time.time() - start_time
        time_record["get_size_bytes"] = sum(
            len(cloudpickle.dumps(obj)) for obj in
            [dag_id, processed_content, question, use_online_model, model_folder,
             temperature, max_token, top_p, repetition_penalty]
        )
        # Build prompt
        # model
        prompt= (
            "#Background#\n"
            "You are a general AI assistant. I will ask you a question. Report your concise thinking thoughts and don't think too complicated, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER].\n"
            "YOUR FINAL ANSWER should be a number OR as few words as possible OR a comma separated list of numbers and/or strings.\n"
            "If you are asked for a number, don’t use comma to write your number neither use units such as $ or percent sign unless specified otherwise.\n"
            "If you are asked for a string, don’t use articles, neither abbreviations (e.g. for cities), and write the digits in plain text unless specified otherwise.\n"
            "If you are asked for a comma separated list, apply the above rules depending of whether the element to be put in the list is a number or a string.\n"
            f"#Question#\n{question}\n"
            "#Extracted text from file#\n"
            f"{processed_content[:8000]}\n"
        )
        payload= {
            "model": "deepseek-ai/DeepSeek-R1-Distill-Qwen-32B",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": temperature,
            "max_tokens": max_token,
            "response_format": {"type": "text"}
        }
        inference_features, answer= None, None
        if use_online_model:
            inference_features, answer = query_llm_online(
                RemoteLlmRoute.from_dag_context(context),
                payload=payload,
                tokenizer_path=tokenizer_path,
            )
        elif backend == "vllm":
            inference_features, answer = query_vllm_model(
                RemoteLlmRoute.for_vllm_worker_base(
                    ray.get(context.get.remote(f"task3_llm_process_deepseek_request_api_url"))
                ),
                model_alias="deepseek-r1-32b",
                messages= payload["messages"],
                temperature= temperature,
                max_token= max_token,
                top_p= top_p,
                repetition_penalty= repetition_penalty)
        else:
            inference_features, answer= query_llm(model_folder= model_folder, model= "deepseek-ai/DeepSeek-R1-Distill-Qwen-32B", messages= [{"role": "user", "content": prompt}], temperature= temperature, max_token= max_token, top_p= top_p, repetition_penalty= repetition_penalty)        
        
        context.put.remote("deepseek_answer", answer)
        inference_features['reason'] = 1
        next_task_prompt = (
            "You are a senior editor and a world-class reasoning expert. Your job is to synthesize the answers from two different AI assistants to produce one final, superior answer for the given question.\n\n"
            f"--- Original Question ---\n{question}\n\n"
            f"--- Answer from Assistant 1 (Qwen3) ---"
            f"--- Answer from Assistant 2 (DeepSeek) ---"
            "--- Your Task ---\n"
            "Analyze both answers. Identify the strengths and weaknesses of each. Then, combine their best elements, correct any errors, and provide a single, comprehensive, and accurate final answer. Adhere to the final answer format requested in the original prompt.\n\n"
            "Report your thoughts, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER]."
        )
        text2_length= len(answer)
        text2_token_count= estimate_tokens(answer)
        text2_feature= {"text2_length": text2_length, "text2_token_count": text2_token_count}
        context.put.remote("text2_feature", text2_feature)
        time_record["put_size_bytes"] += len(cloudpickle.dumps(text2_feature))
        return json.dumps({
            "dag_id": dag_id,
            "curr_task_feat": inference_features,
            "succ_task_feat": {
                "task4_llm_fuse_answer":{
                    "prompt_length": len(next_task_prompt),
                    "prompt_token_count": estimate_tokens(next_task_prompt),         
                    "text2_length": text2_length,
                    "text2_token_count": text2_token_count,
                    "reason": 0
                }
            },            
            "start_time": start_time,
            "end_time": time.time(),
            "time_record": time_record
        })
    except Exception as e:
        print(f"task3_llm_process_deepseek error: {str(e)}")
        raise e

@gpu(gpu_mem=70000, model_name="qwen3-32b", backend="huggingface")
def task4_llm_fuse_answer(context):
    """
    Fusion Task: Takes answers from multiple experts, analyzes them, and generates a final, synthesized answer.
    """
    try:
        #

        time_record = {
            "get_time": 0.0,
            "put_size_bytes": 0,
            "get_size_bytes": 0
        }

        # --- 1.  ---
        start_time = time.time()

        backend = task4_llm_fuse_answer._task_decorator['backend']
        answer_qwen = ray.get(context.get.remote("qwen_answer"))
        answer_deepseek = ray.get(context.get.remote("deepseek_answer"))
        dag_id = ray.get(context.get.remote("dag_id"))
        question = ray.get(context.get.remote("question"))
        text1_feature = ray.get(context.get.remote("text1_feature"))
        text2_feature = ray.get(context.get.remote("text2_feature"))
        use_online_model = ray.get(context.get.remote("use_online_model"))
        model_folder = ray.get(context.get.remote("model_folder"))

        temperature, max_token, repetition_penalty, top_p = None, None, None, None
        if not use_online_model:
            temperature = ray.get(context.get.remote("temperature"))
            max_token = ray.get(context.get.remote("max_tokens"))
            top_p = ray.get(context.get.remote("top_p"))
            repetition_penalty = ray.get(context.get.remote("repetition_penalty"))
            
        hosted_route = None
        vllm_base = None
        if use_online_model:
            hosted_route = RemoteLlmRoute.from_dag_context(context)
        elif backend == "vllm":
            vllm_base = ray.get(context.get.remote(f"task4_llm_fuse_answer_request_api_url"))
        route_meta = (
            hosted_route.redacted_log_hint() if hosted_route else None,
            RemoteLlmRoute.for_vllm_worker_base(vllm_base).redacted_log_hint() if vllm_base else None,
        )

        time_record["get_time"] = time.time() - start_time
        
        # --- 2.  ---
        time_record["get_size_bytes"] = sum(
            len(cloudpickle.dumps(obj)) for obj in
            [answer_qwen, answer_deepseek, dag_id, question, text1_feature, text2_feature,
             use_online_model, model_folder, temperature, max_token, top_p, repetition_penalty,
             route_meta]
        )

        # ---  ---
        print("✅ Task 4 (Fusion): Starting answer synthesis...")
        
        if not question:
            raise ValueError(f"task {dag_id} missing Question")

        tokenizer_path = os.path.join(model_folder, "Qwen/Qwen3-32B")
        prompt = (
            "You are a senior editor and a world-class reasoning expert. Your job is to synthesize the answers from two different AI assistants to produce one final, superior answer for the given question.\n\n"
            f"--- Original Question ---\n{question}\n\n"
            f"--- Answer from Assistant 1 (Qwen3) ---\n{answer_qwen}\n\n"
            f"--- Answer from Assistant 2 (DeepSeek) ---\n{answer_deepseek}\n\n"
            "--- Your Task ---\n"
            "Analyze both answers. Identify the strengths and weaknesses of each. Then, combine their best elements, correct any errors, and provide a single, comprehensive, and accurate final answer. Adhere to the final answer format requested in the original prompt.\n\n"
            "Report your thoughts, and finish your answer with the following template: FINAL ANSWER: [YOUR FINAL ANSWER]."
        )
        
        payload = {
            "model": "Qwen/Qwen3-32B",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": temperature,
            "max_tokens": max_token,
            "enable_thinking": False,
            "response_format": {"type": "text"}
        }

        final_answer = None
        if use_online_model:
            inference_features, final_answer = query_llm_online(
                hosted_route,
                payload=payload,
                tokenizer_path=tokenizer_path,
            )
        elif backend == "vllm":
            inference_features, final_answer = query_vllm_model(
                RemoteLlmRoute.for_vllm_worker_base(vllm_base),
                model_alias="qwen3-32b",
                messages=payload["messages"],
                temperature=temperature,
                max_token=max_token,
                top_p=top_p,
                repetition_penalty=repetition_penalty)
        else:
            inference_features, final_answer = query_llm(model_folder=model_folder, model="Qwen/Qwen3-32B", messages=[{"role": "user", "content": prompt}], temperature=temperature, max_token=max_token, top_p=top_p, repetition_penalty=repetition_penalty)
        
        print(f"  -> Fusion complete. Final answer generated for DAG {dag_id}.")

        # --- 3. ( put put_size_bytes  0) ---
        
        # --- 4.  time_record ---
        return json.dumps({
            "dag_id": dag_id,
            "final_answer": final_answer,
            "curr_task_feat":
            {
                "prompt_length": inference_features["text_length"],
                "prompt_token_count": inference_features["token_count"],
                "text1_length": text1_feature["text1_length"],
                "text1_token_count": text1_feature["text1_token_count"],
                "text2_length": text2_feature["text2_length"],
                "text2_token_count": text2_feature["text2_token_count"],
                "reason": 0
            },
            "start_time": start_time,
            "end_time": time.time(),
            "time_record": time_record #

        })

    except Exception as e:
        print(f"❌ task4_fuse_final_answer failed: {str(e)}")
        raise e